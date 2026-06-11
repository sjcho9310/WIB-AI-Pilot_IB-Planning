"""IB 영업현황 관리 시스템 — PT v2 (reference style)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ── 슬라이드 크기 16:9 ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── 색상 ──────────────────────────────────────────────────────
BG       = RGBColor(0xED, 0xED, 0xED)
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
INK      = RGBColor(0x11, 0x11, 0x11)
GRAY     = RGBColor(0x6B, 0x72, 0x80)
MUTED    = RGBColor(0x9C, 0xA3, 0xAF)
BLUE     = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_LT  = RGBColor(0xEB, 0xF2, 0xFF)
BLUE_DK  = RGBColor(0x1E, 0x3A, 0x8A)
GREEN    = RGBColor(0x05, 0x96, 0x69)
GREEN_LT = RGBColor(0xEC, 0xFF, 0xF4)
RED_FG   = RGBColor(0xDC, 0x26, 0x26)
RED_LT   = RGBColor(0xFE, 0xE2, 0xE2)
LINE     = RGBColor(0xD1, 0xD5, 0xDB)

FONT     = "맑은 고딕"

# ── 공통 헬퍼 ────────────────────────────────────────────────
def bg(sl):
    r = sl.shapes.add_shape(1, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()

def card(sl, x, y, w, h, fill=CARD, border=LINE):
    r = sl.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if border:
        r.line.color.rgb = border; r.line.width = Pt(0.5)
    else:
        r.line.fill.background()
    try: r.adjustments[0] = 0.04
    except: pass
    return r

def txt(sl, text, x, y, w, h, size=Pt(11), bold=False,
        color=INK, align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return tb

def txt_lines(sl, lines, x, y, w, h, size=Pt(11),
              color=GRAY, gap=Pt(3)):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (t, b, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        r = p.add_run(); r.text = t
        r.font.size = size; r.font.bold = b
        r.font.color.rgb = c or color; r.font.name = FONT

def header(sl, chapter, num, total):
    # 상단 가느다란 구분선
    r = sl.shapes.add_shape(1, 0, 0, W, Inches(0.42))
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    txt(sl, chapter, Inches(0.5), Inches(0.1), Inches(6), Inches(0.28),
        size=Pt(10), color=MUTED)
    txt(sl, "IB 영업현황 관리 시스템",
        W-Inches(3.2), Inches(0.1), Inches(2.8), Inches(0.28),
        size=Pt(10), bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

def headline(sl, text, y=Inches(0.55)):
    tb = sl.shapes.add_textbox(Inches(0.5), y, W - Inches(1.0), Inches(1.5))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]
    r  = p.add_run(); r.text = text
    r.font.size = Pt(36); r.font.bold = True
    r.font.color.rgb = INK; r.font.name = FONT

def subhead(sl, text, y=Inches(1.52)):
    txt(sl, text, Inches(0.5), y, W - Inches(1.0), Inches(0.36),
        size=Pt(13), color=GRAY)

def footer(sl, page, source=""):
    txt(sl, page, Inches(0.5), H - Inches(0.3), Inches(1.5), Inches(0.28),
        size=Pt(9), color=MUTED)
    if source:
        txt(sl, source, W - Inches(7.0), H - Inches(0.3), Inches(6.5),
            Inches(0.28), size=Pt(8), color=MUTED, align=PP_ALIGN.RIGHT)

def stat_card(sl, x, y, w, h, num, label, sub="", num_color=BLUE, bg_fill=CARD):
    card(sl, x, y, w, h, fill=bg_fill)
    txt(sl, num,   x+Inches(0.22), y+Inches(0.18), w-Inches(0.3), Inches(0.65),
        size=Pt(38), bold=True, color=num_color)
    txt(sl, label, x+Inches(0.22), y+Inches(0.82), w-Inches(0.3), Inches(0.28),
        size=Pt(11), bold=True, color=INK)
    if sub:
        txt(sl, sub, x+Inches(0.22), y+Inches(1.08), w-Inches(0.3), Inches(0.3),
            size=Pt(10), color=GRAY)

def bullet_card(sl, x, y, w, h, title, items, bg_fill=CARD, title_color=INK):
    card(sl, x, y, w, h, fill=bg_fill)
    txt(sl, title, x+Inches(0.22), y+Inches(0.18), w-Inches(0.3), Inches(0.28),
        size=Pt(13), bold=True, color=title_color)
    tb = sl.shapes.add_textbox(x+Inches(0.22), y+Inches(0.52), w-Inches(0.35),
                               h-Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run(); r1.text = "· "
        r1.font.size = Pt(10.5); r1.font.bold = True
        r1.font.color.rgb = BLUE; r1.font.name = FONT
        r2 = p.add_run(); r2.text = item
        r2.font.size = Pt(10.5); r2.font.color.rgb = GRAY; r2.font.name = FONT

def bar_row(sl, x, y, w, label, pct, val):
    lw = Inches(1.1)
    txt(sl, label, x, y, lw, Inches(0.24),
        size=Pt(10), color=GRAY, align=PP_ALIGN.RIGHT)
    track_x = x + lw + Inches(0.08)
    track_w = w - lw - Inches(0.5)
    r = sl.shapes.add_shape(1, track_x, y, track_w, Inches(0.22))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0xE5,0xE7,0xEB)
    r.line.fill.background()
    if pct > 0:
        fw = max(int(track_w * pct), Inches(0.05))
        f = sl.shapes.add_shape(1, track_x, y, fw, Inches(0.22))
        f.fill.solid(); f.fill.fore_color.rgb = BLUE
        f.line.fill.background()
    txt(sl, val, track_x+track_w+Inches(0.08), y,
        Inches(0.35), Inches(0.24), size=Pt(9), color=MUTED)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — 개발 배경 및 목적
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header(sl, "00  개발 배경 및 목적", "1", "7")
headline(sl, "Excel 분산 관리의 한계,\nIB 딜 파이프라인 통합 디지털화")
subhead(sl, "부서별 수기 입력 → 실시간 통합 플랫폼 · FTP 기반 수익성 자동 산출 · 조직 계층 집계")

# 4개 pain point stat boxes (상단 행)
pain = [
    ("Excel\n분산 관리", "부서별 파일 → 통합 불가", RED_FG, RED_LT),
    ("수기 계산\n오류 빈발", "FTP 원가·수익성 매번 수동 계산", RED_FG, RED_LT),
    ("방치 딜\n추적 불가", "마지막 수정일 관리 수단 없음", BLUE, BLUE_LT),
    ("통합 통계\n부재", "본부·부문 집계 → 보고 지연", BLUE, BLUE_LT),
]
pw = Inches(2.95)
py = Inches(2.05)
ph = Inches(1.38)
for i, (num, sub, fc, bgc) in enumerate(pain):
    px = Inches(0.5) + i * (pw + Inches(0.12))
    card(sl, px, py, pw, ph, fill=bgc, border=None)
    txt(sl, num, px+Inches(0.2), py+Inches(0.1), pw-Inches(0.3), Inches(0.68),
        size=Pt(20), bold=True, color=fc)
    txt(sl, sub, px+Inches(0.2), py+Inches(0.76), pw-Inches(0.3), Inches(0.48),
        size=Pt(10.5), color=GRAY)

# 3개 목적 카드 (하단 행)
purpose = [
    ("딜 파이프라인 통합 관리",
     ["신규딜 입력 → 진행단계 추적 → 인라인 편집",
      "단일 · 복합 Deal 다중 트랜치 구조 지원",
      "조직 권한별 데이터 범위 자동 필터"]),
    ("FTP 기반 수익성 자동 산출",
     ["FTP 선형보간으로 자본원가 자동 계산",
      "Carry손익 = 이자수익 − 자본원가",
      "저장 버튼 클릭 시 전체 재산출"]),
    ("조직 계층 집계 통계",
     ["부서 → 본부 → IB부문 권한별 집계",
      "KPI 5개 · 가로차트 · 요약테이블",
      "경영진 보고용 즉시 출력"]),
]
cw3 = Inches(4.1)
cy3 = Inches(3.6)
ch3 = Inches(3.4)
for i, (ttl, items) in enumerate(purpose):
    cx3 = Inches(0.5) + i * (cw3 + Inches(0.15))
    bullet_card(sl, cx3, cy3, cw3, ch3, ttl, items, title_color=BLUE_DK)

footer(sl, "01 / 07")


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — 화면01 : 조직 선택
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header(sl, "01  조직 선택", "2", "7")
headline(sl, "3단계 선택만으로\n권한과 데이터 범위 자동 결정")
subhead(sl, "로그인 없이 역할 · 본부 · 부서 선택 → session_state에 view_level 저장 → 이후 모든 화면에 자동 반영")

# Flow steps
steps = [
    ("STEP 1", "역할 선택",  "영업조직 / 관리자",           True),
    ("STEP 2", "본부 선택",  "기업금융 / CM / 대체투자\n투자금융 / IB직속", False),
    ("STEP 3", "부서 선택",  "본부 소속 부서\n(영업조직만 해당)",           False),
    ("완료",   "홈 진입",    "view_level 확정",              True),
]
sw = Inches(2.8)
sy = Inches(2.1)
sh = Inches(0.95)
for i, (num, lbl, sub, ac) in enumerate(steps):
    sx = Inches(0.5) + i * (sw + Inches(0.46))
    c_fill = BLUE_LT if ac else CARD
    card(sl, sx, sy, sw, sh, fill=c_fill, border=BLUE if ac else LINE)
    txt(sl, num, sx+Inches(0.15), sy+Inches(0.08), sw, Inches(0.2),
        size=Pt(8.5), bold=True, color=BLUE)
    txt(sl, lbl, sx+Inches(0.15), sy+Inches(0.28), sw-Inches(0.2), Inches(0.28),
        size=Pt(14), bold=True, color=INK)
    txt(sl, sub, sx+Inches(0.15), sy+Inches(0.58), sw-Inches(0.2), Inches(0.35),
        size=Pt(10), color=GRAY)
    if i < 3:
        txt(sl, "→", sx+sw+Inches(0.08), sy+Inches(0.3), Inches(0.35), Inches(0.4),
            size=Pt(18), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# view_level 카드 3개
vl_data = [
    ("department", "영업조직 부서 선택",
     "본인 부서 딜만 조회 · 입력 · 수정 가능", CARD),
    ("division",   "관리자 본부 선택",
     "해당 본부 전 부서 딜 집계 조회", CARD),
    ("all",        "관리자 IB부문통합 선택",
     "전 본부 · 부서 딜 통합 조회 및 통계", BLUE_LT),
]
vw = Inches(3.95)
vy = Inches(3.28)
vh = Inches(2.85)
for i, (tag, ttl, desc, bgc) in enumerate(vl_data):
    vx = Inches(0.5) + i * (vw + Inches(0.22))
    card(sl, vx, vy, vw, vh, fill=bgc, border=BLUE if bgc==BLUE_LT else LINE)
    # tag pill
    pill_w = Inches(1.35)
    p_fill = BLUE if bgc == BLUE_LT else RGBColor(0xE5,0xE7,0xEB)
    p_fc   = CARD if bgc == BLUE_LT else GRAY
    r2 = sl.shapes.add_shape(1, vx+Inches(0.2), vy+Inches(0.2), pill_w, Inches(0.28))
    r2.fill.solid(); r2.fill.fore_color.rgb = p_fill; r2.line.fill.background()
    txt(sl, tag, vx+Inches(0.2), vy+Inches(0.21), pill_w, Inches(0.26),
        size=Pt(9), bold=True, color=p_fc, align=PP_ALIGN.CENTER)
    txt(sl, ttl, vx+Inches(0.2), vy+Inches(0.62), vw-Inches(0.35), Inches(0.32),
        size=Pt(15), bold=True, color=INK)
    txt(sl, desc, vx+Inches(0.2), vy+Inches(1.0), vw-Inches(0.35), Inches(0.5),
        size=Pt(11.5), color=GRAY)
    # divider
    dr = sl.shapes.add_shape(1, vx+Inches(0.2), vy+Inches(1.6), vw-Inches(0.4), Pt(1))
    dr.fill.solid(); dr.fill.fore_color.rgb = LINE; dr.line.fill.background()

    role_txt = "영업조직" if i == 0 else "관리자"
    role_col = BLUE if i > 0 else GREEN
    txt(sl, f"대상: {role_txt}", vx+Inches(0.2), vy+Inches(1.72),
        vw-Inches(0.3), Inches(0.28), size=Pt(10.5), bold=True, color=role_col)

footer(sl, "02 / 07")


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — 화면02 : 홈
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header(sl, "02  홈", "3", "7")
headline(sl, "역할에 따라 자동 활성화되는\n4가지 핵심 기능")
subhead(sl, "우상단 '← 뒤로 / ⌂ 조직선택' 버튼이 모든 화면에 고정 표시 · 조직선택 버튼은 session_state 전체 초기화")

menus = [
    ("✏️", "신규딜 입력",
     ["단일·복합 Deal 구조 분기",
      "3단 연동 드롭다운 입력 폼",
      "당사투자금액=0 시 수익률 비활성"],
     "영업조직 전용", GREEN, GREEN_LT, GREEN),
    ("📋", "딜 리스트",
     ["전체 딜 현황 인라인 편집",
      "FTP 원가·Carry손익 자동 산출",
      "마지막수정일 자동 갱신"],
     "공통", BLUE, BLUE_LT, BLUE),
    ("📊", "영업통계",
     ["KPI 5개 + 가로 막대차트",
      "진행단계·투자유형·사업유형 분석",
      "본부별·부서별 요약테이블"],
     "공통", BLUE, BLUE_LT, BLUE),
    ("⚙️", "FTP 관리",
     ["종금Book / IB Book 금리 입력",
      "9개 만기 구간 선형보간",
      "변경 즉시 전체 수익성 재산출"],
     "관리자 전용", BLUE_DK, BLUE_LT, BLUE_DK),
]
mw = Inches(5.85)
mh = Inches(2.85)
for i, (icon, ttl, items, tag, fc, bgc, tcol) in enumerate(menus):
    mx = Inches(0.5)  + (i % 2) * (mw + Inches(0.45))
    my = Inches(2.12) + (i // 2) * (mh + Inches(0.15))
    card(sl, mx, my, mw, mh, fill=bgc, border=LINE)
    # left accent bar
    ab = sl.shapes.add_shape(1, mx, my, Inches(0.045), mh)
    ab.fill.solid(); ab.fill.fore_color.rgb = fc; ab.line.fill.background()
    # icon + title
    txt(sl, icon, mx+Inches(0.18), my+Inches(0.2), Inches(0.5), Inches(0.4),
        size=Pt(20))
    txt(sl, ttl, mx+Inches(0.7), my+Inches(0.22), Inches(3.2), Inches(0.35),
        size=Pt(16), bold=True, color=INK)
    # tag pill
    pill_w = Inches(1.1)
    pr2 = sl.shapes.add_shape(1, mx+mw-pill_w-Inches(0.18), my+Inches(0.22),
                               pill_w, Inches(0.27))
    pr2.fill.solid(); pr2.fill.fore_color.rgb = fc; pr2.line.fill.background()
    txt(sl, tag, mx+mw-pill_w-Inches(0.18), my+Inches(0.22),
        pill_w, Inches(0.27), size=Pt(8.5), bold=True,
        color=CARD, align=PP_ALIGN.CENTER)
    # bullet items
    for j, item in enumerate(items):
        iy = my + Inches(0.72) + j * Inches(0.6)
        txt(sl, "· " + item, mx+Inches(0.18), iy, mw-Inches(0.3), Inches(0.5),
            size=Pt(11.5), color=GRAY)

footer(sl, "03 / 07")


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — 화면03 : 신규딜 입력
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header(sl, "03  신규딜 입력", "4", "7")
headline(sl, "투자유형 → 상품유형 → Book유형\n3단 연동 드롭다운으로 입력 오류 차단")
subhead(sl, "상위 선택 변경 시 하위 드롭다운 자동 리셋 · 당사투자금액=0이면 수익률 입력 비활성 · 복합 Deal은 트랜치 멀티셀렉트")

# 왼쪽: STEP A / B 구조
LX = Inches(0.5)
LW = Inches(5.8)

card(sl, LX, Inches(2.12), LW, Inches(4.85), fill=CARD)
txt(sl, "입력 구조", LX+Inches(0.25), Inches(2.25), LW, Inches(0.28),
    size=Pt(11), bold=True, color=MUTED)

# STEP A
a_y = Inches(2.6)
r_a = sl.shapes.add_shape(1, LX+Inches(0.2), a_y, LW-Inches(0.4), Inches(1.5))
r_a.fill.solid(); r_a.fill.fore_color.rgb = RGBColor(0xF9,0xFA,0xFB)
r_a.line.color.rgb = LINE; r_a.line.width = Pt(0.5)
txt(sl, "STEP A  —  공통 딜 정보", LX+Inches(0.38), a_y+Inches(0.1),
    LW, Inches(0.25), size=Pt(10), bold=True, color=BLUE)
step_a = "딜명  ·  진행단계  ·  사업유형  ·  투자구조\nCIB여부  ·  모험자본여부  ·  전체딜규모  ·  기표예정일\n상세내용 / 진행상황"
txt(sl, step_a, LX+Inches(0.38), a_y+Inches(0.38), LW-Inches(0.55), Inches(1.0),
    size=Pt(11), color=GRAY)

# STEP B — single / complex branch
b_y = Inches(4.28)
for i, (lbl, col, desc) in enumerate([
    ("단일 Deal", BLUE, "트랜치 섹션 1개\n투자유형·상품유형·Book유형\n수익률·수수료·기간 입력"),
    ("복합 Deal", MUTED,  "트랜치 멀티셀렉트 (2개↑)\n선택값이 각 섹션 헤더로 바인딩\n트랜치 타입 자동 고정"),
]):
    bx = LX + Inches(0.2) + i * (Inches(2.7) + Inches(0.1))
    r_b = sl.shapes.add_shape(1, bx, b_y, Inches(2.7), Inches(2.55))
    r_b.fill.solid(); r_b.fill.fore_color.rgb = (BLUE_LT if i == 0 else RGBColor(0xF9,0xFA,0xFB))
    r_b.line.color.rgb = BLUE if i == 0 else LINE
    r_b.line.width = Pt(0.5)
    # top stripe
    stripe = sl.shapes.add_shape(1, bx, b_y, Inches(2.7), Inches(0.06))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = col; stripe.line.fill.background()
    txt(sl, lbl, bx+Inches(0.15), b_y+Inches(0.12), Inches(2.5), Inches(0.28),
        size=Pt(12), bold=True, color=INK)
    txt(sl, desc, bx+Inches(0.15), b_y+Inches(0.46), Inches(2.42), Inches(2.0),
        size=Pt(11), color=GRAY)

# 오른쪽: cascade + rules
RX = Inches(6.85)
RW = Inches(6.0)

card(sl, RX, Inches(2.12), RW, Inches(4.85), fill=CARD)
txt(sl, "3단 연동 로직", RX+Inches(0.25), Inches(2.25), RW, Inches(0.28),
    size=Pt(11), bold=True, color=MUTED)

# Cascade dark block
csc_y = Inches(2.6)
r_csc = sl.shapes.add_shape(1, RX+Inches(0.2), csc_y, RW-Inches(0.4), Inches(1.88))
r_csc.fill.solid(); r_csc.fill.fore_color.rgb = RGBColor(0x1A,0x1F,0x2E)
r_csc.line.fill.background()

cascade = [
    (Inches(0.12), "투자유형", "→  허용 상품유형 목록 즉시 교체"),
    (Inches(0.42), "상품유형", "→  허용 Book유형 + 수익률구분 교체"),
    (Inches(0.68), "Book유형", "→  FTP 적용 여부 자동 결정"),
]
for indent, tag_t, desc_t in cascade:
    cy_row = csc_y + Inches(0.15) + cascade.index((indent, tag_t, desc_t)) * Inches(0.53)
    pill = sl.shapes.add_shape(1, RX+Inches(0.2)+indent, cy_row,
                               Inches(0.95), Inches(0.3))
    pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(0x1D,0x4E,0xD8)
    pill.line.fill.background()
    txt(sl, tag_t, RX+Inches(0.2)+indent, cy_row+Inches(0.03), Inches(0.95), Inches(0.26),
        size=Pt(9), bold=True, color=CARD, align=PP_ALIGN.CENTER)
    txt(sl, desc_t, RX+Inches(0.2)+indent+Inches(1.05), cy_row+Inches(0.05),
        RW-Inches(0.6)-indent, Inches(0.26), size=Pt(10), color=RGBColor(0xA8,0xB4,0xD0))

# 특수 규칙 3개
rules = [
    ("① 수익률 입력 비활성화",
     "당사투자금액 = 0이면 투자수익률구분·투자수익률 자동 비활성\n→ 주선·인수 전용 딜 오입력 방지"),
    ("② 복합 Deal 트랜치 멀티셀렉트",
     "[단일순위 / 선순위 / 중후순위 / EBL / Mezzanine / Equity]\n선택 순서대로 트랜치 1·2·3 배정, 타입 필드 자동 고정"),
    ("③ 저장 시 deal_id 공유",
     "복합 Deal은 같은 deal_id를 가진 N개 행으로 CSV 저장\n→ 통계에서 deal_id 기준으로 딜 수 집계"),
]
for i, (rtl, rdesc) in enumerate(rules):
    ry2 = Inches(4.6) + i * Inches(0.82)
    r_rule = sl.shapes.add_shape(1, RX+Inches(0.2), ry2, RW-Inches(0.4), Inches(0.72))
    r_rule.fill.solid(); r_rule.fill.fore_color.rgb = RGBColor(0xF9,0xFA,0xFB)
    r_rule.line.color.rgb = LINE; r_rule.line.width = Pt(0.5)
    txt(sl, rtl,   RX+Inches(0.35), ry2+Inches(0.07), RW, Inches(0.22),
        size=Pt(10), bold=True, color=BLUE_DK)
    txt(sl, rdesc, RX+Inches(0.35), ry2+Inches(0.29), RW-Inches(0.5), Inches(0.38),
        size=Pt(10), color=GRAY)

footer(sl, "04 / 07")


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — 화면04 : 딜 리스트
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header(sl, "04  딜 리스트", "5", "7")
headline(sl, "변경된 행만 감지해\n마지막수정일 자동 추적")
subhead(sl, "data_editor 기반 인라인 편집 · 저장 시 FTP 원가·Carry손익·순영업수익 자동 재산출 · 방치 딜 추적 가능")

# 왼쪽: 컬럼 그룹 카드
LX = Inches(0.5)
LW = Inches(6.1)
card(sl, LX, Inches(2.12), LW, Inches(4.85), fill=CARD)
txt(sl, "컬럼 구성", LX+Inches(0.25), Inches(2.25), LW, Inches(0.28),
    size=Pt(11), bold=True, color=MUTED)

col_groups = [
    ("식별 정보",    "최초작성일 · 마지막수정일 · 부서 · 진행단계 · 딜명",     "읽기전용",   MUTED,   RGBColor(0xF3,0xF4,0xF6)),
    ("딜 분류",      "사업유형 · 상품유형 · 투자유형 · 트랜치#",                 "편집가능",   BLUE,    CARD),
    ("금액",         "전체딜규모 · 당사주선인수 · 리소스Book유형 · 투자금액",    "수익 연동",  BLUE,    BLUE_LT),
    ("수익률",       "투자수익률구분 · 투자수익률 · 투자기간만기",                "수익 연동",  BLUE,    BLUE_LT),
    ("자동 산출",    "FTP유형 · FTP원가율 · Carry손익 · 순영업수익",              "읽기전용",   GREEN,   GREEN_LT),
    ("수수료",       "선취수수료구분 · 선취수수료금액",                            "편집가능",   BLUE,    CARD),
]
row_h = Inches(0.64)
for i, (grp, cols, tag_t, tag_c, bgc) in enumerate(col_groups):
    ry3 = Inches(2.58) + i * (row_h + Inches(0.04))
    r3 = sl.shapes.add_shape(1, LX+Inches(0.2), ry3, LW-Inches(0.4), row_h)
    r3.fill.solid(); r3.fill.fore_color.rgb = bgc
    r3.line.color.rgb = LINE; r3.line.width = Pt(0.4)
    txt(sl, grp, LX+Inches(0.35), ry3+Inches(0.08), Inches(1.0), Inches(0.26),
        size=Pt(10), bold=True, color=INK)
    txt(sl, cols, LX+Inches(0.35), ry3+Inches(0.35), Inches(3.6), Inches(0.26),
        size=Pt(9.5), color=GRAY)
    # tag
    tpw = Inches(0.8)
    tp = sl.shapes.add_shape(1, LX+LW-Inches(0.6)-tpw, ry3+Inches(0.18),
                              tpw, Inches(0.25))
    tp.fill.solid(); tp.fill.fore_color.rgb = tag_c; tp.line.fill.background()
    txt(sl, tag_t, LX+LW-Inches(0.6)-tpw, ry3+Inches(0.18), tpw, Inches(0.25),
        size=Pt(7.5), bold=True, color=CARD, align=PP_ALIGN.CENTER)

# 오른쪽: 자동산출 로직 카드 3개
RX = Inches(7.1)
RW = Inches(5.75)
card(sl, RX, Inches(2.12), RW, Inches(4.85), fill=CARD)
txt(sl, "자동 산출 로직", RX+Inches(0.25), Inches(2.25), RW, Inches(0.28),
    size=Pt(11), bold=True, color=MUTED)

logics = [
    ("FTP 원가율",
     "Book 종류에 따라 종금Book / IB Book FTP 테이블 선택\nnumpy.interp()로 투자기간 만기 구간 선형보간 자동 계산\n채무보증Book · Book미사용 · 그룹펀드편입 → 미적용"),
    ("Carry 손익",
     "보유월수 = min(투자기간만기,  기표일→12/31 잔여월수)\nCarry손익 = 투자금액×(수익률−FTP%) / 100 × 보유월수/12\n순영업수익 = 선취수수료 + Carry손익"),
    ("마지막수정일",
     "저장 버튼 클릭 시 deal_id + 트랜치번호 복합키로 비교\n실제 변경된 행만 입력일 컬럼을 오늘 날짜로 업데이트\n→ 방치 딜 추적 · 영업 활동 모니터링 가능"),
]
for i, (ttl3, desc3) in enumerate(logics):
    ly3 = Inches(2.58) + i * Inches(1.46)
    r_l = sl.shapes.add_shape(1, RX+Inches(0.2), ly3, RW-Inches(0.4), Inches(1.32))
    r_l.fill.solid(); r_l.fill.fore_color.rgb = RGBColor(0xF9,0xFA,0xFB)
    r_l.line.color.rgb = BLUE; r_l.line.width = Pt(1.0)
    # top accent
    ta = sl.shapes.add_shape(1, RX+Inches(0.2), ly3, Inches(0.045), Inches(1.32))
    ta.fill.solid(); ta.fill.fore_color.rgb = BLUE; ta.line.fill.background()
    txt(sl, ttl3, RX+Inches(0.38), ly3+Inches(0.1), RW-Inches(0.55), Inches(0.26),
        size=Pt(12), bold=True, color=BLUE_DK)
    txt(sl, desc3, RX+Inches(0.38), ly3+Inches(0.4), RW-Inches(0.55), Inches(0.85),
        size=Pt(10.5), color=GRAY)

footer(sl, "05 / 07")


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — 화면05 : 영업통계
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header(sl, "05  영업통계", "6", "7")
headline(sl, "선택한 조직 범위,\n딜 수·금액·수익 즉시 집계")
subhead(sl, "\"[조직명] 영업통계\" 타이틀 자동 표시 · view_level별 차트 및 요약테이블 자동 분기 · 0값 항목도 항상 표시")

# 5개 KPI
kpis = [
    ("12건",    "전체 딜 수",    "deal_id 기준",      False),
    ("2,450억", "당사 주선금액", "당사주선/인수규모",  False),
    ("1,200억", "당사 투자금액", "정수 단위",         False),
    ("18.5억",  "선취수수료",    "소수점 1자리",      False),
    ("42.3억",  "순영업수익",    "소수점 1자리",      True),
]
kw = Inches(2.38)
ky = Inches(2.1)
kh = Inches(1.32)
for i, (val, lbl, sub, ac) in enumerate(kpis):
    kx = Inches(0.5) + i * (kw + Inches(0.1))
    fc2 = BLUE if ac else BLUE_DK
    bg2 = BLUE_LT if ac else CARD
    card(sl, kx, ky, kw, kh, fill=bg2, border=BLUE if ac else LINE)
    txt(sl, val, kx+Inches(0.18), ky+Inches(0.1), kw-Inches(0.25), Inches(0.65),
        size=Pt(30), bold=True, color=fc2)
    txt(sl, lbl, kx+Inches(0.18), ky+Inches(0.72), kw-Inches(0.25), Inches(0.28),
        size=Pt(11), bold=True, color=INK)
    txt(sl, sub, kx+Inches(0.18), ky+Inches(0.98), kw-Inches(0.25), Inches(0.25),
        size=Pt(9.5), color=GRAY)

# 왼쪽: 차트 2개
LX = Inches(0.5)
LW = Inches(6.1)
card(sl, LX, Inches(3.6), LW, Inches(3.45), fill=CARD)

txt(sl, "진행단계별 딜 수 (정수 축 · 0 기준 고정)", LX+Inches(0.25), Inches(3.72),
    LW, Inches(0.25), size=Pt(9.5), bold=True, color=GRAY)
chart_rows1 = [("내부검토",0.40,"4"), ("IB실무협의",0.30,"3"),
               ("투심위",0.20,"2"), ("승인완료",0.30,"3")]
for i, (lbl, pct, val) in enumerate(chart_rows1):
    bar_row(sl, LX+Inches(0.15), Inches(4.05)+i*Inches(0.38),
            LW-Inches(0.25), lbl, pct, val)

txt(sl, "투자유형별 당사 투자금액 (억원)", LX+Inches(0.25), Inches(5.65),
    LW, Inches(0.25), size=Pt(9.5), bold=True, color=GRAY)
chart_rows2 = [("직접대출",0.70,"840"), ("PI",0.30,"360"),
               ("단순주선",0.0,"0"), ("인수주선",0.10,"120")]
for i, (lbl, pct, val) in enumerate(chart_rows2):
    bar_row(sl, LX+Inches(0.15), Inches(5.98)+i*Inches(0.35),
            LW-Inches(0.25), lbl, pct, val)

# 오른쪽: 요약테이블
RX = Inches(7.1)
RW = Inches(5.75)
card(sl, RX, Inches(3.6), RW, Inches(3.45), fill=CARD)
txt(sl, "진행단계별 요약 테이블", RX+Inches(0.25), Inches(3.72),
    RW, Inches(0.25), size=Pt(9.5), bold=True, color=GRAY)

tbl_hdrs = ["진행단계", "딜수", "주선/인수", "투자금액", "수수료", "순영업수익"]
tbl_rows3 = [
    ("내부검토",   "4", "620억",  "300억", "4.2억", "10.8억"),
    ("IB실무협의", "3", "450억",  "250억", "3.1억",  "8.3억"),
    ("투심위",     "2", "380억",  "180억", "2.5억",  "6.2억"),
    ("승인완료",   "3", "890억",  "450억", "7.8억", "14.5억"),
]
thws3 = [Inches(1.1), Inches(0.42), Inches(0.88), Inches(0.88), Inches(0.72), Inches(1.0)]
trh3 = Inches(0.34)
for ri, row_d in enumerate([tbl_hdrs]+list(tbl_rows3)):
    for ci, (cell, thw3) in enumerate(zip(row_d, thws3)):
        rx5 = RX+Inches(0.18)+sum(thws3[:ci])
        ry5 = Inches(4.02)+ri*trh3
        is_h = ri == 0
        bg5 = RGBColor(0xF3,0xF4,0xF6) if is_h else (RGBColor(0xF9,0xFA,0xFB) if ri%2 else CARD)
        r5 = sl.shapes.add_shape(1, rx5, ry5, thw3, trh3)
        r5.fill.solid(); r5.fill.fore_color.rgb = bg5
        r5.line.color.rgb = LINE; r5.line.width = Pt(0.3)
        txt(sl, cell, rx5+Inches(0.05), ry5+Inches(0.06),
            thw3-Inches(0.06), trh3-Inches(0.08),
            size=Pt(9), bold=is_h, color=MUTED if is_h else INK,
            align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)

# view_level 표시 note
note_y = Inches(5.82)
r_note = sl.shapes.add_shape(1, RX+Inches(0.18), note_y, RW-Inches(0.3), Inches(0.95))
r_note.fill.solid(); r_note.fill.fore_color.rgb = BLUE_LT
r_note.line.color.rgb = BLUE; r_note.line.width = Pt(0.5)
txt_lines(sl, [
    ("view_level별 추가 표시", True, BLUE_DK),
    ("division → 부서별 요약 테이블 추가", False, GRAY),
    ("all (통합) → 본부별 요약 + 부서별 차트 추가", False, GRAY),
], RX+Inches(0.3), note_y+Inches(0.07), RW-Inches(0.45), Inches(0.85),
   size=Pt(10), gap=Pt(2))

footer(sl, "06 / 07")


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — 화면06 : FTP 관리
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
header(sl, "06  FTP 관리", "7", "7")
headline(sl, "9개 만기 구간 금리 입력,\n전체 딜 수익성 즉시 재산출")
subhead(sl, "관리자 전용 · 종금Book / IB Book 구분 입력 · 저장 즉시 모든 딜의 Carry손익 · 순영업수익 자동 업데이트")

# 왼쪽: FTP 입력 UI mockup
LX = Inches(0.5)
LW = Inches(6.1)
card(sl, LX, Inches(2.12), LW, Inches(4.85), fill=RGBColor(0x16,0x18,0x1C), border=None)

# 탭 바
r_tab = sl.shapes.add_shape(1, LX+Inches(0.18), Inches(2.3), LW-Inches(0.35), Inches(0.4))
r_tab.fill.solid(); r_tab.fill.fore_color.rgb = RGBColor(0x1E,0x22,0x28)
r_tab.line.fill.background()
r_act = sl.shapes.add_shape(1, LX+Inches(0.22), Inches(2.32), Inches(2.55), Inches(0.35))
r_act.fill.solid(); r_act.fill.fore_color.rgb = BLUE; r_act.line.fill.background()
txt(sl, "종금Book FTP", LX+Inches(0.22), Inches(2.34), Inches(2.55), Inches(0.3),
    size=Pt(9.5), bold=True, color=CARD, align=PP_ALIGN.CENTER)
txt(sl, "IB Book FTP", LX+Inches(2.88), Inches(2.34), Inches(2.55), Inches(0.3),
    size=Pt(9.5), color=RGBColor(0x6B,0x72,0x80), align=PP_ALIGN.CENTER)

txt(sl, "9개 만기 구간 입력 (% 단위)", LX+Inches(0.25), Inches(2.82),
    LW, Inches(0.24), size=Pt(9), color=RGBColor(0x6B,0x72,0x80))

tenors = [("1d","3.85"),("1m","3.90"),("3m","3.95"),
          ("6m","4.00"),("1y","4.10"),("2y","4.20"),
          ("3y","4.35"),("5y","4.50"),("10y","4.70")]
tw2 = Inches(1.82)
th2 = Inches(0.62)
for i, (tenor, rate) in enumerate(tenors):
    cx2 = LX + Inches(0.2) + (i%3)*(tw2+Inches(0.08))
    cy2 = Inches(3.12) + (i//3)*(th2+Inches(0.08))
    r_t = sl.shapes.add_shape(1, cx2, cy2, tw2, th2)
    r_t.fill.solid(); r_t.fill.fore_color.rgb = RGBColor(0x1E,0x22,0x2A)
    r_t.line.color.rgb = RGBColor(0x2C,0x30,0x3A); r_t.line.width = Pt(0.5)
    txt(sl, tenor, cx2+Inches(0.12), cy2+Inches(0.06), tw2, Inches(0.22),
        size=Pt(9), color=RGBColor(0x6B,0x72,0x80))
    txt(sl, rate, cx2+Inches(0.12), cy2+Inches(0.3), tw2, Inches(0.28),
        size=Pt(16), bold=True, color=RGBColor(0x93,0xC5,0xFD))

r_save = sl.shapes.add_shape(1, LX+Inches(0.2), Inches(5.84), LW-Inches(0.38), Inches(0.42))
r_save.fill.solid(); r_save.fill.fore_color.rgb = BLUE; r_save.line.fill.background()
txt(sl, "저장", LX+Inches(0.2), Inches(5.87), LW-Inches(0.38), Inches(0.35),
    size=Pt(11), bold=True, color=CARD, align=PP_ALIGN.CENTER)

# 오른쪽: book 적용 기준 + 보간
RX = Inches(7.1)
RW = Inches(5.75)
card(sl, RX, Inches(2.12), RW, Inches(2.38), fill=CARD)
txt(sl, "Book유형별 FTP 적용 기준", RX+Inches(0.25), Inches(2.22),
    RW, Inches(0.25), size=Pt(11), bold=True, color=MUTED)

book_rows2 = [
    ("종금Book",                    "적용", "종금Book FTP"),
    ("셀다운 · PI · ECM DCM · 메자닌", "적용", "IB Book FTP"),
    ("채무보증Book",                 "미적용", "자본원가 = 0"),
    ("Book미사용 · 그룹펀드편입",    "미적용", "자본원가 = 0"),
]
bthws2 = [Inches(2.5), Inches(0.72), Inches(2.08)]
brh2 = Inches(0.35)
for ri, (book2, apply2, tbl2) in enumerate([("Book유형","FTP 적용","사용 테이블")] + book_rows2):
    is_h4 = ri == 0
    for ci2, (val2, thw4) in enumerate(zip([book2,apply2,tbl2], bthws2)):
        rx7 = RX+Inches(0.2)+sum(bthws2[:ci2])
        ry7 = Inches(2.55)+ri*brh2
        bg7 = RGBColor(0xF3,0xF4,0xF6) if is_h4 else (CARD if ri%2 else RGBColor(0xF9,0xFA,0xFB))
        r7 = sl.shapes.add_shape(1, rx7, ry7, thw4, brh2)
        r7.fill.solid(); r7.fill.fore_color.rgb = bg7
        r7.line.color.rgb = LINE; r7.line.width = Pt(0.3)
        if ci2 == 1 and not is_h4:
            ac_b = apply2 == "적용"
            bp = sl.shapes.add_shape(1, rx7+Inches(0.05), ry7+Inches(0.06),
                                     Inches(0.58), Inches(0.22))
            bp.fill.solid()
            bp.fill.fore_color.rgb = BLUE_LT if ac_b else RED_LT
            bp.line.fill.background()
            txt(sl, apply2, rx7+Inches(0.05), ry7+Inches(0.07), Inches(0.58), Inches(0.2),
                size=Pt(8), bold=True,
                color=BLUE if ac_b else RED_FG, align=PP_ALIGN.CENTER)
        else:
            txt(sl, val2, rx7+Inches(0.08), ry7+Inches(0.07), thw4-Inches(0.1), brh2-Inches(0.1),
                size=Pt(9), bold=is_h4, color=MUTED if is_h4 else INK)

# 선형 보간 카드
interp_y = Inches(4.62)
card(sl, RX, interp_y, RW, Inches(1.18), fill=CARD)
ib2 = sl.shapes.add_shape(1, RX, interp_y, Inches(0.045), Inches(1.18))
ib2.fill.solid(); ib2.fill.fore_color.rgb = BLUE; ib2.line.fill.background()
txt(sl, "선형 보간  (Linear Interpolation)", RX+Inches(0.2), interp_y+Inches(0.1),
    RW, Inches(0.26), size=Pt(11), bold=True, color=BLUE_DK)
formula_bg = sl.shapes.add_shape(1, RX+Inches(0.2), interp_y+Inches(0.42),
                                  RW-Inches(0.35), Inches(0.65))
formula_bg.fill.solid()
formula_bg.fill.fore_color.rgb = RGBColor(0xF1,0xF5,0xFD)
formula_bg.line.fill.background()
txt(sl, "numpy.interp(18, [12, 24], [4.10, 4.20])  →  4.15%\n(1y~2y 사이 18개월: 비례 보간)",
    RX+Inches(0.3), interp_y+Inches(0.45), RW-Inches(0.45), Inches(0.55),
    size=Pt(10.5), color=BLUE_DK)

# 전후 변화
chg_y = Inches(5.92)
card(sl, RX, chg_y, RW, Inches(1.0), fill=BLUE_LT, border=BLUE)
txt(sl, "FTP 미등록 시  →  FTP원가율 · Carry손익 · 순영업수익 모두 \"–\" 표시",
    RX+Inches(0.22), chg_y+Inches(0.1), RW-Inches(0.3), Inches(0.3),
    size=Pt(10.5), color=BLUE_DK)
txt(sl, "FTP 등록 후  →  전체 딜 자동 재산출 · FTP 변경 이력 누적 보관",
    RX+Inches(0.22), chg_y+Inches(0.42), RW-Inches(0.3), Inches(0.3),
    size=Pt(10.5), color=BLUE_DK)

footer(sl, "07 / 07",
       "Streamlit PoC · Python 3.11 · CSV 저장소 · FTP 선형보간 · Firebase Studio")

# ── 저장 ────────────────────────────────────────────────────
OUT = "/home/user/test_repository/presentation/IB_영업현황_관리시스템_PT.pptx"
prs.save(OUT)
print(f"저장 완료: {OUT}")
